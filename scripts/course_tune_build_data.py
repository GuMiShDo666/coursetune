#!/usr/bin/env python3
"""Build source chunks, annotation prompts, and training data for a course assistant."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from pathlib import Path


SUPPORTED_SUFFIXES = {".pdf", ".pptx", ".txt", ".md"}
SYSTEM_PROMPT = "你是 EBU5606 产品开发课程资料智能答疑助手。只根据课程资料回答；如果资料不足，就说明无法从资料中确定。"


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", nargs="+", required=True, help="Course files or folders.")
    parser.add_argument("--out", default="data/course_assistant_annotation_prompts.jsonl")
    parser.add_argument(
        "--mode",
        choices=["chunks", "sft_prompts", "dpo_prompts", "sft_dataset", "dpo_dataset"],
        default="sft_prompts",
    )
    parser.add_argument("--max-chars", type=int, default=1400)
    parser.add_argument("--limit", type=int, default=0)
    parser.add_argument("--name-regex", default="", help="Optional regex filter applied to each file name.")
    parser.add_argument("--seed-json", default="", help="Optional reviewed JSON sample file prepended to training data.")
    parser.add_argument("--seed-repeat", type=int, default=1, help="Number of times to repeat seed rows.")
    args = parser.parse_args()

    name_pattern = re.compile(args.name_regex) if args.name_regex else None
    files = list(iter_files([Path(item).expanduser() for item in args.source], name_pattern=name_pattern))
    if not files:
        print("No supported files found. Supported suffixes: .pdf, .pptx, .txt, .md", file=sys.stderr)
        return 2

    rows = []
    for file_path in files:
        for section in extract_file(file_path):
            for chunk_index, chunk in enumerate(chunk_text(section["text"], max_chars=args.max_chars), start=1):
                chunk_row = {
                    "id": build_id(file_path, section["locator"], chunk_index),
                    "source_file": str(file_path.resolve()),
                    "source_name": file_path.name,
                    "locator": section["locator"],
                    "text": chunk,
                }
                rows.extend(convert_rows(args.mode, chunk_row))
                if args.limit and len(rows) >= args.limit:
                    break
            if args.limit and len(rows) >= args.limit:
                break
        if args.limit and len(rows) >= args.limit:
            break

    seed_rows = load_seed_rows(args.seed_json, args.mode, repeat=args.seed_repeat)
    if seed_rows:
        rows = seed_rows + rows

    write_jsonl(args.out, rows)
    print(f"Wrote {len(rows)} rows to {args.out}")
    return 0


def iter_files(paths: list[Path], name_pattern: re.Pattern[str] | None = None):
    seen = set()
    for path in paths:
        if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES and matches_name(path, name_pattern):
            resolved = path.resolve()
            if resolved not in seen:
                seen.add(resolved)
                yield resolved
        elif path.is_dir():
            for child in sorted(path.rglob("*")):
                if child.is_file() and child.suffix.lower() in SUPPORTED_SUFFIXES and matches_name(child, name_pattern):
                    resolved = child.resolve()
                    if resolved not in seen:
                        seen.add(resolved)
                        yield resolved


def matches_name(path: Path, name_pattern: re.Pattern[str] | None) -> bool:
    return name_pattern is None or bool(name_pattern.search(path.name))


def extract_file(file_path: Path) -> list[dict[str, str]]:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(file_path)
    if suffix == ".pptx":
        return extract_pptx(file_path)
    return [{"locator": "file", "text": file_path.read_text(encoding="utf-8", errors="ignore")}]


def extract_pdf(file_path: Path) -> list[dict[str, str]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Install PDF support with: pip install -r requirements/course_tune.txt") from exc

    reader = PdfReader(str(file_path))
    rows = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            rows.append({"locator": f"page:{page_number}", "text": text})
    return rows


def extract_pptx(file_path: Path) -> list[dict[str, str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Install PPTX support with: pip install -r requirements/course_tune.txt") from exc

    presentation = Presentation(str(file_path))
    rows = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        if parts:
            rows.append({"locator": f"slide:{slide_number}", "text": "\n".join(parts)})
    return rows


def convert_rows(mode: str, chunk: dict[str, str]) -> list[dict[str, object]]:
    if mode == "chunks":
        return [chunk]

    if mode == "sft_prompts":
        prompt = f"""你是数据标注员。请根据下面课程资料生成 3 条中文 SFT 训练样本。

要求：
1. 只能使用资料中的信息，不要补充外部知识。
2. 问题要覆盖定义、比较、机制或应用。
3. 输出 JSONL，每行包含 instruction, input, output, system, source_files。

来源：{chunk["source_name"]} {chunk["locator"]}

资料：
{chunk["text"]}"""
        return [annotation_prompt_row(mode, chunk, prompt)]

    if mode == "dpo_prompts":
        prompt = f"""你是偏好数据标注员。请根据下面课程资料生成 2 条中文 DPO 偏好样本。

要求：
1. chosen 必须准确、清楚、忠于资料。
2. rejected 应包含常见错误、遗漏关键条件或答非所问。
3. 输出 JSONL，每行包含 instruction, input, chosen, rejected, system, source_files。

来源：{chunk["source_name"]} {chunk["locator"]}

资料：
{chunk["text"]}"""
        return [annotation_prompt_row(mode, chunk, prompt)]

    if mode == "sft_dataset":
        return build_sft_rows(chunk)

    return build_dpo_rows(chunk)


def annotation_prompt_row(mode: str, chunk: dict[str, str], prompt: str) -> dict[str, str]:
    return {
        "id": f"{mode}_{chunk['id']}",
        "source_file": chunk["source_file"],
        "locator": chunk["locator"],
        "prompt": prompt,
    }


def build_sft_rows(chunk: dict[str, str]) -> list[dict[str, object]]:
    topic = topic_hint(chunk["text"], chunk["source_name"])
    answer = format_course_answer(chunk)
    source = f"{chunk['source_name']} {chunk['locator']}"
    return [
        {
            "instruction": f"根据 EBU5606 课程资料，解释「{topic}」这部分内容。",
            "input": "",
            "output": answer,
            "system": SYSTEM_PROMPT,
            "source_files": [chunk["source_name"]],
        },
        {
            "instruction": "请把下面这段 EBU5606 课程资料整理成中文复习笔记。",
            "input": f"来源：{source}\n\n资料：\n{chunk['text']}",
            "output": answer,
            "system": SYSTEM_PROMPT,
            "source_files": [chunk["source_name"]],
        },
    ]


def build_dpo_rows(chunk: dict[str, str]) -> list[dict[str, object]]:
    topic = topic_hint(chunk["text"], chunk["source_name"])
    chosen = format_course_answer(chunk)
    rejected = (
        f"「{topic}」主要是一个通用管理概念，回答时可以结合常识自由发挥。"
        "不需要严格区分课件中的阶段、条件或来源，也不需要说明资料是否不足。"
    )
    return [
        {
            "instruction": f"根据 EBU5606 课程资料，解释「{topic}」这部分内容。",
            "input": "",
            "chosen": chosen,
            "rejected": rejected,
            "system": SYSTEM_PROMPT,
            "source_files": [chunk["source_name"]],
        }
    ]


def format_course_answer(chunk: dict[str, str]) -> str:
    bullets = build_bullets(chunk["text"])
    if bullets:
        body = "\n".join(f"- {item}" for item in bullets)
    else:
        body = chunk["text"]
    return f"根据课程资料，这部分的核心内容是：\n{body}\n\n来源：{chunk['source_name']}，{chunk['locator']}。"


def build_bullets(text: str, limit: int = 5) -> list[str]:
    candidates = []
    for line in text.splitlines():
        line = line.strip(" -•\t")
        if len(line) < 8:
            continue
        candidates.append(line)
        if len(candidates) >= limit:
            break
    if not candidates and text.strip():
        candidates = [text.strip()[:500]]
    return candidates


def topic_hint(text: str, source_name: str, max_chars: int = 48) -> str:
    for line in text.splitlines():
        cleaned = re.sub(r"\s+", " ", line).strip(" -•\t")
        if len(cleaned) >= 4 and not is_generic_heading(cleaned):
            return cleaned[:max_chars]
    source_topic = re.sub(r"\.[^.]+$", "", source_name)
    source_topic = re.sub(r"^EBU5606\s*[-_]\s*", "", source_topic, flags=re.IGNORECASE)
    source_topic = re.sub(r"_?2026$", "", source_topic)
    source_topic = re.sub(r"\s+", " ", source_topic).strip(" -_")
    return source_topic[:max_chars] or "课程知识点"


def is_generic_heading(text: str) -> bool:
    normalized = text.lower().strip()
    generic = {
        "ebu5606",
        "product development",
        "product development and marketing",
        "queen mary university of london",
    }
    return normalized in generic or normalized.startswith("slide ")


def chunk_text(text: str, max_chars: int) -> list[str]:
    normalized = normalize_text(text)
    if not normalized:
        return []
    if len(normalized) <= max_chars:
        return [normalized]
    chunks = []
    start = 0
    while start < len(normalized):
        chunk = normalized[start : start + max_chars].strip()
        if chunk:
            chunks.append(chunk)
        start += max_chars
    return chunks


def normalize_text(text: str) -> str:
    lines = []
    for line in text.replace("\x00", " ").splitlines():
        compact = re.sub(r"\s+", " ", line).strip()
        if compact:
            lines.append(compact)
    return "\n".join(lines)


def build_id(file_path: Path, locator: str, chunk_index: int) -> str:
    raw = f"{file_path.resolve()}::{locator}::{chunk_index}"
    return hashlib.sha1(raw.encode("utf-8")).hexdigest()[:16]


def load_seed_rows(path: str, mode: str, repeat: int) -> list[dict[str, object]]:
    if not path or mode not in {"sft_dataset", "dpo_dataset"}:
        return []
    seed_path = Path(path).expanduser()
    if not seed_path.exists():
        raise FileNotFoundError(f"Seed JSON not found: {seed_path}")
    rows = json.loads(seed_path.read_text(encoding="utf-8"))
    if not isinstance(rows, list):
        raise ValueError("Seed JSON must contain a list of training rows.")

    required = {"instruction", "input", "system"}
    if mode == "sft_dataset":
        required.add("output")
    else:
        required.update({"chosen", "rejected"})

    cleaned = []
    for index, row in enumerate(rows, start=1):
        if not isinstance(row, dict) or not required.issubset(row):
            raise ValueError(f"Seed row {index} does not match {mode} columns.")
        cleaned.append(row)
    return cleaned * max(1, repeat)


def write_jsonl(path: str, rows: list[dict[str, object]]) -> None:
    out = Path(path)
    out.parent.mkdir(parents=True, exist_ok=True)
    if out.suffix.lower() == ".json":
        with out.open("w", encoding="utf-8") as handle:
            json.dump(rows, handle, ensure_ascii=False, indent=2)
            handle.write("\n")
        return
    with out.open("w", encoding="utf-8") as handle:
        for row in rows:
            handle.write(json.dumps(row, ensure_ascii=False))
            handle.write("\n")


if __name__ == "__main__":
    raise SystemExit(main())
